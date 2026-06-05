"""Condense the 3-slide GenBI x Hilton deck onto ONE slide.

Three columns — Solution · Architecture · Get Started — plus a closing strip.

Run:  python build_genbi_onepager.py
Out:  GenBI_Hilton_OnePager.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette ---------------------------------------------------------------------
INK      = RGBColor(0x0F, 0x17, 0x2A)
SLATE    = RGBColor(0x33, 0x41, 0x55)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
BG       = RGBColor(0xF8, 0xFA, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_L = RGBColor(0x63, 0x66, 0xF1)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
CYAN_D   = RGBColor(0x06, 0xB6, 0xD4)
CARD_TINT= RGBColor(0xEE, 0xF2, 0xFF)
SF_BLUE  = RGBColor(0x03, 0x2D, 0x60)
SF_SKY   = RGBColor(0x00, 0xA1, 0xE0)
SLATE300 = RGBColor(0x94, 0xA3, 0xB8)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG


def shape(x, y, w, h, fill=None, line=None, line_w=1.0,
          kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def txt(x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if "space_after" in p: para.space_after = Pt(p["space_after"])
        if "line" in p: para.line_spacing = p["line"]
        r = para.add_run(); r.text = p["t"]
        f = r.font; f.size = Pt(p.get("size", 12)); f.bold = p.get("bold", False)
        f.name = FONT; f.color.rgb = p.get("color", INK)
    return tb


def ctext(sp, lines, size=11, color=WHITE, bold=True):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [{"t": lines, "size": size, "color": color, "bold": bold}]
    for i, p in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        if "space_after" in p: para.space_after = Pt(p["space_after"])
        r = para.add_run(); r.text = p["t"]
        r.font.size = Pt(p.get("size", size)); r.font.bold = p.get("bold", bold)
        r.font.name = FONT; r.font.color.rgb = p.get("color", color)
    return sp


def chip(x, y, w, text_, fill=CARD_TINT, color=INDIGO, size=9, h=0.3):
    sp = shape(x, y, w, h, fill=fill, radius=0.5)
    ctext(sp, text_, size=size, color=color)
    sp.text_frame.word_wrap = False
    return sp


def down_arrow(cx, y, w=0.34, h=0.26, color=INDIGO_L):
    return shape(cx - w / 2, y, w, h, fill=color, kind=MSO_SHAPE.DOWN_ARROW)


def bullet(x, y, w, text_, dot=INDIGO, size=9.5):
    shape(x, y + 0.04, 0.13, 0.13, fill=dot, kind=MSO_SHAPE.OVAL)
    txt(x + 0.24, y - 0.03, w - 0.24, 0.45, [{"t": text_, "size": size, "color": SLATE,
        "line": 1.1}])


# ============================================================== HEADER ========
shape(0, 0, 13.333, 0.16, fill=INDIGO, kind=MSO_SHAPE.RECTANGLE)
txt(0.5, 0.32, 3.0, 0.5, [{"t": "Gen", "size": 20, "bold": True, "color": INK}],
    anchor=MSO_ANCHOR.MIDDLE)
txt(0.5 + 0.56, 0.32, 3.0, 0.5, [{"t": "BI", "size": 20, "bold": True, "color": CYAN_D}],
    anchor=MSO_ANCHOR.MIDDLE)
txt(2.2, 0.30, 10.6, 0.55, [{"t": "Agentic, self-serve BI — native to Salesforce",
    "size": 21, "bold": True, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
txt(0.5, 0.92, 12.3, 0.4, [{"t": "Non-technical Hilton users turn fragmented data into "
    "board-ready dashboards in minutes — inside the Salesforce they already use.",
    "size": 12, "color": SLATE}])

# ============================================================== COLUMNS ========
MARGIN = 0.5
GAP = 0.34
CW = (13.333 - 2 * MARGIN - 2 * GAP) / 3        # column width
TOP = 1.55
CARDH = 5.05
cols = [MARGIN + i * (CW + GAP) for i in range(3)]
HEADS = [("1 · THE SOLUTION", "BI for everyone — no SQL", INDIGO),
         ("2 · ARCHITECTURE", "An LWC inside Salesforce", SF_BLUE),
         ("3 · GET STARTED", "Live in Hilton's SF org", CYAN_D)]

for cx, (h1, h2, accent) in zip(cols, HEADS):
    shape(cx, TOP, CW, CARDH, fill=WHITE, line=LINE, line_w=1, radius=0.045)
    shape(cx, TOP, CW, 0.62, fill=accent, radius=0.045)
    shape(cx, TOP + 0.31, CW, 0.31, fill=accent, kind=MSO_SHAPE.RECTANGLE)
    txt(cx + 0.2, TOP + 0.08, CW - 0.4, 0.5, [
        {"t": h1, "size": 11.5, "bold": True, "color": WHITE, "space_after": 1},
        {"t": h2, "size": 9.5, "color": RGBColor(0xE2, 0xE8, 0xF0)}])

ix = lambda c: c + 0.22                          # inner-x for a column
iw = CW - 0.44                                    # inner width
body = TOP + 0.78                                 # content start y

# ---- COLUMN 1 : Solution ----------------------------------------------------
c = cols[0]
ctext(shape(ix(c), body, iw, 0.4, fill=CARD_TINT, radius=0.18),
      "Fragmented data", size=10.5, color=INDIGO)
down_arrow(c + CW / 2, body + 0.45, color=INDIGO_L)
ctext(shape(ix(c), body + 0.76, iw, 0.92, fill=INK, radius=0.1), [
    {"t": "AGENTIC PROCESS", "size": 9.5, "color": CYAN, "bold": True, "space_after": 2},
    {"t": "Cleanse · Merge · Classify", "size": 10, "color": WHITE, "space_after": 1},
    {"t": "Analyze → your KPIs", "size": 10, "color": WHITE}])
down_arrow(c + CW / 2, body + 1.74, color=CYAN_D)
ctext(shape(ix(c), body + 2.05, iw, 0.4, fill=CARD_TINT, radius=0.18),
      "Aesthetic, live dashboard", size=10.5, color=CYAN_D)
# feature bullets
fb = body + 2.7
txt(ix(c), fb, iw, 0.3, [{"t": "PLUS", "size": 8.5, "bold": True, "color": MUTED}])
bullet(ix(c), fb + 0.30, iw, "Build & change dashboards just by chatting", dot=INDIGO)
bullet(ix(c), fb + 0.85, iw, "BI Companion watches your work & suggests next steps",
       dot=CYAN_D)

# ---- COLUMN 2 : Architecture ------------------------------------------------
c = cols[1]
# Salesforce container
sfh = 1.95
shape(ix(c), body, iw, sfh, fill=WHITE, line=SF_SKY, line_w=1.25, radius=0.06)
shape(ix(c), body, iw, 0.36, fill=SF_BLUE, radius=0.06)
shape(ix(c), body + 0.18, iw, 0.18, fill=SF_BLUE, kind=MSO_SHAPE.RECTANGLE)
txt(ix(c) + 0.12, body, iw - 0.24, 0.36, [{"t": "SALESFORCE PLATFORM · Hilton org",
    "size": 9.5, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
ctext(shape(ix(c) + 0.16, body + 0.48, iw - 0.32, 0.78, fill=INK, radius=0.1), [
    {"t": "GenBI — embedded LWC", "size": 10.5, "color": CYAN, "bold": True, "space_after": 2},
    {"t": "Dashboards + BI Companion chat", "size": 9, "color": WHITE}])
txt(ix(c) + 0.16, body + 1.34, iw - 0.32, 0.5, [{"t": "Salesforce objects · SSO / OAuth · "
    "Permission sets", "size": 8.5, "color": SF_BLUE, "bold": True, "line": 1.1}])
# connector
down_arrow(c + CW / 2, body + sfh + 0.04, w=0.4, h=0.3, color=INDIGO)
txt(ix(c), body + sfh + 0.04, iw, 0.3, [{"t": "secure API · REST / SSE", "size": 8,
    "bold": True, "color": MUTED, "align": PP_ALIGN.RIGHT}])
# GenBI services container
gy = body + sfh + 0.42
shape(ix(c), gy, iw, 1.18, fill=WHITE, line=INDIGO, line_w=1.25, radius=0.06)
shape(ix(c), gy, iw, 0.34, fill=INDIGO, radius=0.06)
shape(ix(c), gy + 0.17, iw, 0.17, fill=INDIGO, kind=MSO_SHAPE.RECTANGLE)
txt(ix(c) + 0.12, gy, iw - 0.24, 0.34, [{"t": "GenBI AGENTIC SERVICES · cloud",
    "size": 9.5, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
txt(ix(c) + 0.16, gy + 0.44, iw - 0.32, 0.7, [
    {"t": "LangGraph agent engine", "size": 9.5, "bold": True, "color": INK, "space_after": 1},
    {"t": "PostgreSQL store · OpenAI", "size": 9, "color": MUTED}])

# ---- COLUMN 3 : Get Started -------------------------------------------------
c = cols[2]
steps = [
    ("Install the GenBI package / LWC into Hilton's Salesforce org.", INDIGO),
    ("Configure trust — Connected App, Named Credential & Remote Site.", INDIGO_L),
    ("Drop the GenBI LWC onto a Lightning page in App Builder.", VIOLET),
    ("Connect Salesforce + external data and pick Hilton's KPIs.", CYAN_D),
    ("Assign permission sets to business SPOCs — go live.", CYAN),
]
sy = body + 0.02
rowh = 0.84
for i, (t, accent) in enumerate(steps):
    y = sy + i * rowh
    ctext(shape(ix(c), y, 0.46, 0.46, fill=accent, kind=MSO_SHAPE.OVAL),
          str(i + 1), size=15, color=WHITE)
    txt(ix(c) + 0.62, y - 0.02, iw - 0.62, rowh, [{"t": t, "size": 10, "color": SLATE,
        "line": 1.12}], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================== FOOTER =========
ctext(shape(MARGIN, 6.78, 13.333 - 2 * MARGIN, 0.5, fill=INK, radius=0.14),
      "Self-serve, agentic BI for every Hilton SPOC — delivered where they already work, "
      "inside Salesforce.", size=12.5, color=WHITE)

out = "GenBI_Hilton_OnePager.pptx"
prs.save(out)
print(f"Saved {out}")
