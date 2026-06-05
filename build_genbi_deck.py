"""Generate the GenBI x Hilton pitch deck (3 slides), editable .pptx.

Slides:
  1. Solution Overview & Benefits  — agentic BI for non-technical users
  2. Technical Architecture        — GenBI as an LWC inside Salesforce
  3. Getting Started               — quick setup in Hilton's Salesforce org

Run:  python build_genbi_deck.py
Out:  GenBI_Hilton_Overview.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette ----
INK      = RGBColor(0x0F, 0x17, 0x2A)
SLATE    = RGBColor(0x33, 0x41, 0x55)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
BG       = RGBColor(0xF8, 0xFA, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_L = RGBColor(0x63, 0x66, 0xF1)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
CYAN_D   = RGBColor(0x06, 0xB6, 0xD4)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
CARD_TINT= RGBColor(0xEE, 0xF2, 0xFF)
INK_CARD = RGBColor(0x1B, 0x24, 0x40)
SF_BLUE  = RGBColor(0x03, 0x2D, 0x60)   # Salesforce navy
SF_SKY   = RGBColor(0x00, 0xA1, 0xE0)   # Salesforce cloud blue
SLATE300 = RGBColor(0x94, 0xA3, 0xB8)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers ----
def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def shape(s, x, y, w, h, fill=None, line=None, line_w=1.0,
          kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def rect(s, x, y, w, h, **kw):
    return shape(s, x, y, w, h, **kw)


def txt(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if "space_after" in p:  para.space_after = Pt(p["space_after"])
        if "space_before" in p: para.space_before = Pt(p["space_before"])
        if "line" in p:         para.line_spacing = p["line"]
        r = para.add_run(); r.text = p["t"]
        f = r.font
        f.size = Pt(p.get("size", 14)); f.bold = p.get("bold", False)
        f.name = p.get("font", FONT); f.color.rgb = p.get("color", INK)
    return tb


def center_text(sp, text_, size=11, color=WHITE, bold=True, font=FONT, lines=None):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    items = lines if lines else [{"t": text_, "size": size, "color": color, "bold": bold}]
    for i, p in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        if "space_after" in p: para.space_after = Pt(p["space_after"])
        r = para.add_run(); r.text = p["t"]
        r.font.size = Pt(p.get("size", size)); r.font.bold = p.get("bold", bold)
        r.font.name = font; r.font.color.rgb = p.get("color", color)
    return sp


def chip(s, x, y, w, text_, fill=CARD_TINT, color=INDIGO, size=10.5, h=0.34, bold=True):
    sp = rect(s, x, y, w, h, fill=fill, radius=0.5)
    center_text(sp, text_, size=size, color=color, bold=bold)
    sp.text_frame.word_wrap = False
    return sp


def logo(s, x, y, dark=False):
    base = WHITE if dark else INK
    txt(s, x, y, 3.2, 0.5, [{"t": "Gen", "size": 22, "bold": True, "color": base}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.62, y, 3.0, 0.5, [{"t": "BI", "size": 22, "bold": True, "color": CYAN_D}],
        anchor=MSO_ANCHOR.MIDDLE)


def eyebrow(s, x, y, label, num):
    txt(s, x, y, 9, 0.3, [{"t": label.upper(), "size": 12, "bold": True, "color": INDIGO}])
    txt(s, 12.4, 0.42, 0.6, 0.4, [{"t": num, "size": 12, "bold": True, "color": MUTED,
                                    "align": PP_ALIGN.RIGHT}])


def header(s, eyebrow_text, title, num, subtitle=None):
    rect(s, 0, 0, 13.333, 0.16, fill=INDIGO, kind=MSO_SHAPE.RECTANGLE)
    logo(s, 0.6, 0.42)
    eyebrow(s, 0.6, 1.12, eyebrow_text, num)
    txt(s, 0.6, 1.40, 12.1, 0.7, [{"t": title, "size": 25, "bold": True, "color": INK}])
    if subtitle:
        txt(s, 0.6, 2.06, 12.1, 0.6, [{"t": subtitle, "size": 13, "color": SLATE, "line": 1.18}])


def arrow(s, x, y, w=0.45, h=0.5, color=INDIGO_L, kind=MSO_SHAPE.CHEVRON):
    return rect(s, x, y, w, h, fill=color, kind=kind)


def icon_dot(s, x, y, d=0.24, fill=INDIGO):
    return rect(s, x, y, d, d, fill=fill, kind=MSO_SHAPE.OVAL)


# ================================================================ SLIDE 1 ====
def slide1():
    s = add_slide()
    header(s, "Solution Overview & Benefits",
           "Business intelligence for everyone — not just data teams", "01",
           "Non-technical users get board-ready insights in minutes. An agentic workflow "
           "consolidates fragmented data, and an AI Companion guides every step.")

    band_y, band_h = 2.92, 1.92
    # --- Zone A: fragmented data -------------------------------------------
    ax, aw = 0.6, 2.75
    rect(s, ax, band_y, aw, band_h, fill=WHITE, line=LINE, line_w=1)
    txt(s, ax, band_y + 0.12, aw, 0.3, [{"t": "FRAGMENTED DATA", "size": 10, "bold": True,
        "color": MUTED, "align": PP_ALIGN.CENTER}])
    srcs = [("Databases", CYAN_D), ("Spreadsheets", GREEN),
            ("Files / CSV", VIOLET), ("Salesforce objects", SF_SKY)]
    iy = band_y + 0.52
    for name, c in srcs:
        icon_dot(s, ax + 0.22, iy + 0.04, d=0.2, fill=c)
        txt(s, ax + 0.52, iy, aw - 0.6, 0.3, [{"t": name, "size": 10.5, "bold": True,
            "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
        iy += 0.33

    # chevron A->B
    arrow(s, ax + aw + 0.08, band_y + band_h/2 - 0.28, w=0.40, h=0.56, color=INDIGO_L)

    # --- Zone B: agentic process (dark) ------------------------------------
    bx, bw = 4.05, 4.95
    rect(s, bx, band_y, bw, band_h, fill=INK)
    txt(s, bx, band_y + 0.12, bw, 0.3, [{"t": "AGENTIC PROCESS", "size": 10.5, "bold": True,
        "color": CYAN, "align": PP_ALIGN.CENTER}])
    steps = [("Cleanse", INDIGO), ("Merge", INDIGO),
             ("Classify", VIOLET), ("Analyze → your KPIs", VIOLET)]
    pw = (bw - 0.5 - 0.18) / 2; ph = 0.44
    p0x = bx + 0.25
    for i, (st, c) in enumerate(steps):
        col = i % 2; row = i // 2
        pxx = p0x + col * (pw + 0.18)
        pyy = band_y + 0.46 + row * (ph + 0.10)
        center_text(rect(s, pxx, pyy, pw, ph, fill=c, radius=0.22), st, size=11, color=WHITE)
    txt(s, bx, band_y + band_h - 0.30, bw, 0.3, [{"t": "Autonomous AI agents — no SQL, no code",
        "size": 9.5, "color": SLATE300, "align": PP_ALIGN.CENTER}])

    # chevron B->C
    arrow(s, bx + bw + 0.08, band_y + band_h/2 - 0.28, w=0.40, h=0.56, color=CYAN_D)

    # --- Zone C: aesthetic dashboard ---------------------------------------
    cx, cw = 9.55, 3.18
    rect(s, cx, band_y, cw, band_h, fill=WHITE, line=LINE, line_w=1)
    txt(s, cx, band_y + 0.12, cw, 0.3, [{"t": "AESTHETIC DASHBOARD", "size": 10, "bold": True,
        "color": MUTED, "align": PP_ALIGN.CENTER}])
    # KPI tiles
    kt_y = band_y + 0.5
    center_text(rect(s, cx + 0.2, kt_y, 1.34, 0.46, fill=CARD_TINT, radius=0.18),
                "ADR  $182", size=10, color=INDIGO)
    center_text(rect(s, cx + 1.64, kt_y, 1.34, 0.46, fill=CARD_TINT, radius=0.18),
                "RevPAR $146", size=10, color=CYAN_D)
    # mini bar chart
    base = band_y + band_h - 0.30
    heights = [0.45, 0.72, 0.55, 0.95, 0.78]
    bxx = cx + 0.28
    for i, hh in enumerate(heights):
        col = INDIGO if i % 2 == 0 else CYAN_D
        rect(s, bxx, base - hh, 0.20, hh, fill=col, kind=MSO_SHAPE.RECTANGLE)
        bxx += 0.30
    # mini donut
    rect(s, cx + 2.05, band_y + 1.06, 0.78, 0.78, fill=VIOLET, kind=MSO_SHAPE.OVAL)
    rect(s, cx + 2.27, band_y + 1.28, 0.34, 0.34, fill=WHITE, kind=MSO_SHAPE.OVAL)

    # zone captions
    cap_y = band_y + band_h + 0.06
    for zx, zw, cap in [(ax, aw, "Siloed & messy"), (bx, bw, "Consolidated automatically"),
                        (cx, cw, "Live & shareable")]:
        txt(s, zx, cap_y, zw, 0.3, [{"t": cap, "size": 10.5, "bold": True, "color": SLATE,
            "align": PP_ALIGN.CENTER}])

    # --- Feature callouts ---------------------------------------------------
    fy, fh = 5.55, 1.45
    # Card 1: build by chatting
    c1x, c1w = 0.6, 6.0
    rect(s, c1x, fy, c1w, fh, fill=WHITE, line=LINE, line_w=1)
    rect(s, c1x, fy, 0.09, fh, fill=INDIGO, kind=MSO_SHAPE.RECTANGLE)
    txt(s, c1x + 0.28, fy + 0.14, c1w - 0.5, 0.3, [{"t": "Build & change dashboards by chatting",
        "size": 13.5, "bold": True, "color": INK}])
    center_text(rect(s, c1x + 0.28, fy + 0.55, 3.4, 0.4, fill=INDIGO, radius=0.4),
                "“Add ADR by region as a bar chart”", size=10, color=WHITE)
    center_text(rect(s, c1x + 1.9, fy + 1.0, 3.8, 0.38, fill=CARD_TINT, radius=0.4),
                "✓ Added ‘ADR by Region’ — anything else?", size=10, color=INDIGO)

    # Card 2: BI Companion
    c2x, c2w = 6.83, 5.9
    rect(s, c2x, fy, c2w, fh, fill=WHITE, line=LINE, line_w=1)
    rect(s, c2x, fy, 0.09, fh, fill=CYAN_D, kind=MSO_SHAPE.RECTANGLE)
    txt(s, c2x + 0.28, fy + 0.14, c2w - 0.5, 0.3, [{"t": "BI Companion — it watches & suggests",
        "size": 13.5, "bold": True, "color": INK}])
    center_text(rect(s, c2x + 0.28, fy + 0.55, c2w - 0.55, 0.46, fill=INK, radius=0.28),
                "Suggestion: Occupancy dipped in Q2 — see a trend by property?",
                size=10, color=WHITE)
    txt(s, c2x + 0.28, fy + 1.08, c2w - 0.55, 0.3, [{"t": "Context-aware, next-best-analysis "
        "prompts as you work.", "size": 10, "color": MUTED}])


# ================================================================ SLIDE 2 ====
def slide2():
    s = add_slide()
    header(s, "Technical Architecture · Salesforce-native", "01".replace("01", ""), "02")
    # custom title (overwrite)
    txt(s, 0.6, 1.40, 12.1, 0.7, [{"t": "Delivered as a Lightning Web Component inside Salesforce",
        "size": 24, "bold": True, "color": INK}])
    txt(s, 0.6, 2.04, 12.1, 0.5, [{"t": "GenBI runs natively in your Salesforce org — secure, "
        "governed, zero context-switching. The LWC calls the GenBI agentic services over a "
        "trusted connection.", "size": 12.5, "color": SLATE, "line": 1.18}])

    top = 2.72; ch = 3.55
    # --- Salesforce platform container -------------------------------------
    sx, sw = 0.6, 6.85
    rect(s, sx, top, sw, ch, fill=WHITE, line=SF_SKY, line_w=1.5, radius=0.035)
    rect(s, sx, top, sw, 0.5, fill=SF_BLUE, kind=MSO_SHAPE.RECTANGLE)
    txt(s, sx + 0.2, top, sw - 0.4, 0.5, [{"t": "SALESFORCE PLATFORM  (Hilton org)",
        "size": 12.5, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    # Lightning App Builder page
    lx, lw, ly, lh = sx + 0.25, sw - 0.5, top + 0.66, 1.85
    rect(s, lx, ly, lw, lh, fill=BG, line=SF_SKY, line_w=1)
    txt(s, lx + 0.15, ly + 0.08, lw - 0.3, 0.3, [{"t": "Lightning App Builder page  ·  "
        "Home / App / Record", "size": 10.5, "bold": True, "color": SF_BLUE}])
    # GenBI LWC component (prominent)
    gx, gw, gy, gh = lx + 0.25, lw - 0.5, ly + 0.45, 1.18
    rect(s, gx, gy, gw, gh, fill=INK)
    center_text(rect(s, gx, gy, gw, gh, fill=INK, radius=0.08), "", lines=[
        {"t": "GenBI  —  embedded LWC", "size": 13, "color": CYAN, "bold": True, "space_after": 3},
        {"t": "Interactive dashboards  +  BI Companion chat", "size": 10.5, "color": WHITE},
        {"t": "renders inside the page, like any native component", "size": 9, "color": SLATE300},
    ])
    # SF supporting chips
    chip_y = ly + lh + 0.16
    chips = [("Standard & custom objects", SF_BLUE), ("SSO / OAuth identity", SF_BLUE),
             ("Permission sets", SF_BLUE)]
    ccx = sx + 0.25
    for name, c in chips:
        w = 0.35 + 0.085 * len(name)
        chip(s, ccx, chip_y, w, name, fill=RGBColor(0xE6,0xF2,0xFB), color=SF_BLUE, size=10)
        ccx += w + 0.18
    txt(s, sx + 0.25, chip_y + 0.46, sw - 0.5, 0.3, [{"t": "Data, identity & access stay "
        "governed by Salesforce.", "size": 9.5, "color": MUTED}])

    # --- connection arrows --------------------------------------------------
    midx = sx + sw + 0.04
    arrow(s, midx, top + 1.15, w=0.62, h=0.5, color=INDIGO, kind=MSO_SHAPE.RIGHT_ARROW)
    arrow(s, midx, top + 2.05, w=0.62, h=0.5, color=CYAN_D, kind=MSO_SHAPE.LEFT_ARROW)
    txt(s, midx - 0.05, top + 0.78, 0.85, 0.3, [{"t": "secure", "size": 8.5, "bold": True,
        "color": MUTED, "align": PP_ALIGN.CENTER}])
    txt(s, midx - 0.05, top + 2.58, 0.85, 0.3, [{"t": "JSON / SSE", "size": 8.5, "bold": True,
        "color": MUTED, "align": PP_ALIGN.CENTER}])

    # --- GenBI agentic services container ----------------------------------
    ex, ew = 8.2, 4.53
    rect(s, ex, top, ew, ch, fill=WHITE, line=INDIGO, line_w=1.5, radius=0.035)
    rect(s, ex, top, ew, 0.5, fill=INDIGO, kind=MSO_SHAPE.RECTANGLE)
    txt(s, ex + 0.2, top, ew - 0.4, 0.5, [{"t": "GenBI AGENTIC SERVICES  (cloud)",
        "size": 12, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    inner_x = ex + 0.25; inner_w = ew - 0.5
    center_text(rect(s, inner_x, top + 0.66, inner_w, 0.42, fill=CARD_TINT, radius=0.14),
                "API layer — REST + streaming (SSE)", size=10.5, color=INDIGO)
    center_text(rect(s, inner_x, top + 1.18, inner_w, 0.74, fill=INK, radius=0.12), "", lines=[
        {"t": "Agentic engine · LangGraph", "size": 11, "color": CYAN, "bold": True, "space_after": 2},
        {"t": "Cleanse · Merge · Classify · Analyze · KPIs", "size": 9.5, "color": WHITE}])
    center_text(rect(s, inner_x, top + 2.02, inner_w * 0.58 - 0.08, 0.56, fill=BG,
                     line_w=1, line=LINE), "Governed store\nPostgreSQL", lines=[
        {"t": "Governed store", "size": 10, "color": INK, "bold": True, "space_after": 1},
        {"t": "PostgreSQL", "size": 9, "color": MUTED}])
    center_text(rect(s, inner_x + inner_w * 0.58 + 0.08, top + 2.02, inner_w * 0.42 - 0.0, 0.56,
                     fill=BG, line_w=1, line=LINE), "", lines=[
        {"t": "LLM", "size": 10, "color": INK, "bold": True, "space_after": 1},
        {"t": "OpenAI", "size": 9, "color": MUTED}])
    txt(s, inner_x, top + 2.72, inner_w, 0.3, [{"t": "Ingests SF objects + external sources",
        "size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER}])

    # --- footer tech chips --------------------------------------------------
    txt(s, 0.6, 6.42, 6, 0.3, [{"t": "INTEGRATION", "size": 10, "bold": True, "color": MUTED}])
    stack = ["LWC", "Apex callouts", "Named Credentials", "Remote Site", "OAuth 2.0",
             "REST + SSE", "PostgreSQL"]
    cx = 0.6
    for st in stack:
        w = 0.38 + 0.095 * len(st)
        chip(s, cx, 6.74, w, st, fill=CARD_TINT, color=INDIGO, size=10)
        cx += w + 0.16


# ================================================================ SLIDE 3 ====
def slide3():
    s = add_slide()
    header(s, "Getting Started", "Live in Hilton's Salesforce org in five quick steps", "03",
           "No data migration. No new logins. BI delivered where Hilton teams already work.")

    steps = [
        ("Install the package", "Deploy the GenBI managed package / LWC into Hilton's "
         "Salesforce org.", INDIGO),
        ("Establish trust", "Configure the Connected App, Named Credential & Remote Site "
         "to reach GenBI services securely.", INDIGO_L),
        ("Place the component", "In Lightning App Builder, drop the GenBI LWC onto a "
         "Home, App or Record page.", VIOLET),
        ("Connect data & KPIs", "Link Salesforce objects + external sources; pick the KPIs "
         "Hilton cares about.", CYAN_D),
        ("Assign & go live", "Grant permission sets to business SPOCs — they start "
         "chatting with their data.", CYAN),
    ]

    # process ribbon (chevrons) ---------------------------------------------
    rib_y, rib_h = 2.95, 0.66
    margin = 0.6
    total_w = 13.333 - 2 * margin
    n = len(steps)
    overlap = 0.18
    cw = (total_w + (n - 1) * overlap) / n
    for i, (t, d, accent) in enumerate(steps):
        x = margin + i * (cw - overlap)
        kind = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        center_text(rect(s, x, rib_y, cw, rib_h, fill=accent, kind=kind),
                    f"{i+1}.  {t}", size=11.5, color=WHITE)

    # step description cards -------------------------------------------------
    cy, cardh = 3.95, 2.3
    gap = 0.28
    cardw = (total_w - (n - 1) * gap) / n
    for i, (t, d, accent) in enumerate(steps):
        x = margin + i * (cardw + gap)
        rect(s, x, cy, cardw, cardh, fill=WHITE, line=LINE, line_w=1)
        rect(s, x, cy, cardw, 0.12, fill=accent, kind=MSO_SHAPE.RECTANGLE)
        center_text(rect(s, x + (cardw - 0.62) / 2, cy + 0.28, 0.62, 0.62, fill=accent,
                         kind=MSO_SHAPE.OVAL), str(i + 1), size=22, color=WHITE)
        txt(s, x + 0.18, cy + 1.05, cardw - 0.36, cardh - 1.1, [
            {"t": t, "size": 12.5, "bold": True, "color": INK, "align": PP_ALIGN.CENTER,
             "space_after": 6, "line": 1.05},
            {"t": d, "size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER, "line": 1.22},
        ])

    # CTA bar
    center_text(rect(s, 0.6, 6.55, 12.13, 0.6, fill=INK, radius=0.12),
                "Result: every Hilton SPOC gets self-serve, agentic BI — inside Salesforce.",
                size=13, color=WHITE)


slide1()
slide2()
slide3()

out = "GenBI_Hilton_Overview.pptx"
prs.save(out)
print(f"Saved {out}")
